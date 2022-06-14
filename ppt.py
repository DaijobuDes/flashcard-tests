import pptx

# Documentation: https://python-pptx.readthedocs.io/en/latest/

def is_underlined(text) -> False:
    return text.font.underline

def is_bold(text) -> False:
    return text.font.bold

def is_term(text) -> None:
    return is_underlined(text) or is_bold(text)


def main(data: pptx.Presentation) -> None:
    text_runs = []

    for slide in data.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if is_bold(run):
                        # print(f"Term: {run.text}")
                        text_runs.append(run.text)
                    elif is_underlined(run):
                        # print(f"Def: {run.text}")
                        text_runs.append(run.text)

    return text_runs


if __name__ == '__main__':
    data = pptx.Presentation(".\\test.pptx")
    output = main(data)
    print(output)